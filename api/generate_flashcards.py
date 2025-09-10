# api/generate_flashcards.py
import json
import io
import os
import base64
from typing import List
import pdfplumber
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate, HumanMessagePromptTemplate
from langchain_core.output_parsers import PydanticOutputParser
from pydantic import BaseModel, Field
from vercel_event_handler import VercelHandler   # <-- NEW import

# ----- Pydantic models -----
class Flashcard(BaseModel):
    question: str = Field(description="The question on the flashcard")
    answer: str = Field(description="The answer on the flashcard")

class FlashcardList(BaseModel):
    flashcards: List[Flashcard] = Field(description="List of flashcards")

# ----- Text extraction -----
def extract_text_from_pdf(file_content):
    try:
        with pdfplumber.open(io.BytesIO(file_content)) as pdf:
            return "\n".join([page.extract_text() or "" for page in pdf.pages])
    except Exception as e:
        return f"Error reading PDF: {e}"

def extract_text_from_docx(file_content):
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_content))
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    except Exception as e:
        return f"Error reading DOCX: {e}"

def extract_text_from_pptx(file_content):
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return f"Error reading PPTX: {e}"

# ----- Flashcard generation -----
def generate_flashcards(text, api_key):
    parser = PydanticOutputParser(pydantic_object=FlashcardList)
    prompt = ChatPromptTemplate.from_messages([
        HumanMessagePromptTemplate.from_template(
            """You are a flashcard generator for theory-based subjects.
            You must ONLY use information that appears in the provided text.
            Generate as many flashcards as possible (aim for at least 20 if content allows).
            Each flashcard must:
            - Have a clear question
            - Provide a 2-3 sentence answer
            - Stay strictly factual, based only on the provided text
            {format_instructions}
            Text: {input_text}"""
        )
    ]).partial(format_instructions=parser.get_format_instructions())

    llm = ChatOpenAI(
        model="gpt-4o-mini",
        temperature=0.3,
        api_key=api_key,
        max_tokens=2000
    )
    chain = prompt | llm | parser

    # Chunk text to avoid token limits
    chunk_size = 4000
    chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
    all_flashcards = []
    for chunk in chunks:
        if chunk.strip():
            result = chain.invoke({"input_text": chunk})
            all_flashcards.extend(result.flashcards)
    return [{"question": c.question, "answer": c.answer} for c in all_flashcards]

# ----- Main logic -----
def lambda_handler(event, context=None):
    try:
        if event['httpMethod'] == 'OPTIONS':
            return {
                'statusCode': 200,
                'headers': {
                    'Access-Control-Allow-Origin': '*',
                    'Access-Control-Allow-Methods': 'POST, OPTIONS',
                    'Access-Control-Allow-Headers': 'Content-Type',
                },
                'body': ''
            }

        if event['httpMethod'] != 'POST':
            return {
                'statusCode': 405,
                'headers': {'Content-Type': 'application/json'},
                'body': json.dumps({'success': False, 'error': 'Method not allowed'})
            }

        # Get OpenAI API key
        api_key = os.environ.get('OPENAI_API_KEY')
        if not api_key:
            return {
                'statusCode': 500,
                'headers': {'Content-Type': 'application/json'},
                'body': json.dumps({'success': False, 'error': 'OpenAI API key not configured'})
            }

        # Parse JSON body
        content_type = event['headers'].get('content-type', '')
        if 'application/json' not in content_type:
            return {
                'statusCode': 400,
                'headers': {'Content-Type': 'application/json'},
                'body': json.dumps({'success': False, 'error': 'Content-Type must be application/json'})
            }

        body = json.loads(event['body'])
        file_content = base64.b64decode(body['file_content'])
        file_type = body['file_type']

        # Extract text
        if file_type == 'application/pdf':
            text = extract_text_from_pdf(file_content)
        elif file_type == 'text/plain':
            text = file_content.decode('utf-8')
        elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            text = extract_text_from_docx(file_content)
        elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            text = extract_text_from_pptx(file_content)
        else:
            return {
                'statusCode': 400,
                'headers': {'Content-Type': 'application/json'},
                'body': json.dumps({'success': False, 'error': 'Unsupported file type'})
            }

        if not text.strip():
            return {
                'statusCode': 400,
                'headers': {'Content-Type': 'application/json'},
                'body': json.dumps({'success': False, 'error': 'Could not extract text from file'})
            }

        # Generate flashcards
        flashcards = generate_flashcards(text, api_key)
        return {
            'statusCode': 200,
            'headers': {
                'Content-Type': 'application/json',
                'Access-Control-Allow-Origin': '*'
            },
            'body': json.dumps({
                'success': True,
                'flashcards': flashcards,
                'count': len(flashcards)
            })
        }

    except Exception as e:
        # Return JSON even on error
        return {
            'statusCode': 500,
            'headers': {
                'Content-Type': 'application/json',
                'Access-Control-Allow-Origin': '*'
            },
            'body': json.dumps({
                'success': False,
                'error': str(e)
            })
        }

# ----- Export for Vercel -----
handler = VercelHandler(lambda_handler)
